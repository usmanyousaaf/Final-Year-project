from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import tempfile
import os
import requests
import re
import json

app = Flask(__name__)

# Configuration
GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
MODEL_NAME = "gemma2-9b-it"

# Style Options
FONT_CHOICES = {
    "Arial": "Arial",
    "Calibri": "Calibri",
    "Times New Roman": "Times New Roman",
    "Helvetica": "Helvetica",
    "Verdana": "Verdana",
    "Georgia": "Georgia",
    "Courier New": "Courier New",
    "Palatino": "Palatino",
    "Garamond": "Garamond",
    "Trebuchet MS": "Trebuchet MS"
}

COLOR_PALETTES = {
    "Professional": {"primary": "#2B579A", "secondary": "#5B9BD5", "text": "#000000", "background": "#FFFFFF"},
    "Modern": {"primary": "#404040", "secondary": "#A5A5A5", "text": "#FFFFFF", "background": "#121212"},
    "Vibrant": {"primary": "#E53935", "secondary": "#FFCDD2", "text": "#212121", "background": "#F5F5F5"},
    "Corporate": {"primary": "#1976D2", "secondary": "#BBDEFB", "text": "#212121", "background": "#FAFAFA"},
    "Creative": {"primary": "#7B1FA2", "secondary": "#CE93D8", "text": "#FFFFFF", "background": "#212121"},
    "Nature": {"primary": "#388E3C", "secondary": "#A5D6A7", "text": "#212121", "background": "#F1F8E9"},
    "Elegant": {"primary": "#5D4037", "secondary": "#BCAAA4", "text": "#FFFFFF", "background": "#3E2723"},
    "Minimal": {"primary": "#000000", "secondary": "#E0E0E0", "text": "#212121", "background": "#FFFFFF"}
}

def get_ai_response(prompt: str) -> str:
    """Get formatted content from Groq API."""
    if not GROQ_API_KEY:
        return "API key not provided. Please set GROQ_API_KEY in environment variables."
    
    headers = {"Authorization": f"Bearer {GROQ_API_KEY}"}
    payload = {
        "messages": [{
            "role": "user",
            "content": f"{prompt}\n\nReturn well-structured content with clear headings, concise bullet points (max 5 per slide), and no markdown. Use bold/italic via formatting, not symbols."
        }],
        "model": MODEL_NAME,
        "temperature": 0.7
    }
    
    try:
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers=headers,
            json=payload,
            timeout=30
        )
        return response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        return f"Error: {str(e)}"

def extract_slide_text(slide) -> str:
    """Extract text from slide."""
    return "\n".join(
        paragraph.text
        for shape in slide.shapes
        if hasattr(shape, "text_frame")
        for paragraph in shape.text_frame.paragraphs
    )

def split_formatted_runs(text):
    """Split text into formatted runs removing markdown."""
    formatted_segments = []
    bold_parts = re.split(r'(\*\*)', text)
    bold = False
    current_text = []
    
    for part in bold_parts:
        if part == '**':
            if current_text:
                formatted_segments.append({'text': ''.join(current_text), 'bold': bold, 'italic': False})
                current_text = []
            bold = not bold
        else:
            italic_parts = re.split(r'(\*)', part)
            italic = False
            for italic_part in italic_parts:
                if italic_part == '*':
                    if current_text:
                        formatted_segments.append({'text': ''.join(current_text), 'bold': bold, 'italic': italic})
                        current_text = []
                    italic = not italic
                else:
                    current_text.append(italic_part)
            if current_text:
                formatted_segments.append({'text': ''.join(current_text), 'bold': bold, 'italic': italic})
                current_text = []
    return formatted_segments

def apply_formatting(text_frame, content, design_settings):
    """Apply proper formatting to text frame with markdown removal."""
    text_frame.clear()
    paragraphs = [p for p in content.split('\n') if p.strip()]
    
    for para in paragraphs:
        clean_para = re.sub(r'^#+\s*', '', para).strip()
        
        if re.match(r'^(What is|Understanding|What\'s|Drive Theory|Incentive Theory)', clean_para, re.IGNORECASE):
            p = text_frame.add_paragraph()
            p.text = clean_para
            p.font.size = Pt(design_settings["heading_size"])
            p.font.bold = True
            p.space_after = Pt(12)
            continue
            
        if re.match(r'^[\•\-\*] ', clean_para):
            p = text_frame.add_paragraph()
            p.level = 0
            content = clean_para[2:].strip()
            segments = split_formatted_runs(content)
            for seg in segments:
                run = p.add_run()
                run.text = seg['text']
                run.font.bold = seg['bold']
                run.font.italic = seg['italic']
            p.font.size = Pt(design_settings["body_size"])
            continue
            
        p = text_frame.add_paragraph()
        segments = split_formatted_runs(clean_para)
        for seg in segments:
            run = p.add_run()
            run.text = seg['text']
            run.font.bold = seg['bold']
            run.font.italic = seg['italic']
        p.font.size = Pt(design_settings["body_size"])
        p.space_after = Pt(6)

def enhance_slide(slide, user_prompt, design_settings):
    """Enhance slide content with proper formatting."""
    original_text = extract_slide_text(slide)
    if not original_text.strip():
        return
    
    prompt = f"Improve this slide content:\n{original_text}\n\nInstructions: {user_prompt}"
    enhanced_content = get_ai_response(prompt)
    
    title_shape = getattr(slide.shapes, 'title', None)
    shapes_to_remove = [s for s in slide.shapes if s != title_shape]
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    left = Inches(1) if title_shape else Inches(0.5)
    top = Inches(1.5) if title_shape else Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, Inches(8), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    apply_formatting(text_frame, enhanced_content, design_settings)
    apply_design(slide, design_settings)

def apply_design(slide, design_settings):
    """Apply visual design to slide."""
    colors = design_settings["colors"]
    
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = design_settings["font"]
                    run.font.color.rgb = RGBColor.from_string(colors["text"][1:])
                    
                    if shape == getattr(slide.shapes, 'title', None):
                        run.font.color.rgb = RGBColor.from_string(colors["primary"][1:])
                        run.font.size = Pt(design_settings["title_size"])
                        run.font.bold = True
    
    if design_settings["set_background"]:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(colors["background"][1:])

def process_presentation(file_data, user_prompt, design_settings):
    """Process and enhance the entire presentation."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(file_data)
        tmp_path = tmp_file.name
    
    prs = Presentation(tmp_path)
    
    for slide in prs.slides:
        enhance_slide(slide, user_prompt, design_settings)
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    os.unlink(tmp_path)
    return output

@app.route('/')
def index():
    return render_template('index.html', 
                          fonts=FONT_CHOICES, 
                          color_palettes=COLOR_PALETTES)
                          
@app.route('/debug')
def debug():
    """Simple debug page for file uploads"""
    return render_template('debug.html')

@app.route('/enhance', methods=['POST'])
def enhance():
    # Debug information
    app.logger.info("Files received: %s", request.files)
    app.logger.info("Form data: %s", request.form)
    
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if not file.filename.endswith('.pptx'):
        return jsonify({"error": "Only PPTX files are allowed"}), 400
    
    user_prompt = request.form.get('prompt', '')
    
    # Parse design settings from form
    design_settings = {
        "font": request.form.get('font', 'Calibri'),
        "color_palette": request.form.get('colorPalette', 'Professional'),
        "title_size": int(request.form.get('titleSize', 32)),
        "heading_size": int(request.form.get('headingSize', 24)),
        "body_size": int(request.form.get('bodySize', 18)),
        "set_background": request.form.get('setBackground', 'false') == 'true'
    }
    
    # Check if using custom colors
    use_custom = request.form.get('useCustom', 'false') == 'true'
    if use_custom:
        custom_colors = {
            "primary": request.form.get('primaryColor', '#2B579A'),
            "secondary": request.form.get('secondaryColor', '#5B9BD5'),
            "text": request.form.get('textColor', '#000000'),
            "background": request.form.get('backgroundColor', '#FFFFFF')
        }
        design_settings["colors"] = custom_colors
    else:
        design_settings["colors"] = COLOR_PALETTES[design_settings["color_palette"]]
    
    try:
        # Save the file data
        file_data = file.read()
        app.logger.info("File successfully read, size: %d bytes", len(file_data))
        
        # Process the presentation
        output = process_presentation(file_data, user_prompt, design_settings)
        app.logger.info("Presentation processing completed")
        
        # Return the file for download
        return send_file(
            output,
            as_attachment=True,
            download_name="enhanced.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        app.logger.error("Error processing presentation: %s", str(e))
        return jsonify({"error": str(e)}), 500

@app.route('/check_api_key')
def check_api_key():
    if GROQ_API_KEY:
        return jsonify({"status": "API key is set"})
    else:
        return jsonify({"status": "API key is not set"}), 400

@app.route('/upload_debug', methods=['POST'])
def upload_debug():
    """Debug route for file uploads"""
    app.logger.info("Files received: %s", request.files)
    app.logger.info("Form data: %s", request.form)
    
    if 'file' not in request.files:
        return jsonify({"status": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"status": "No selected file"}), 400
    
    filename = file.filename
    file_size = 0
    try:
        data = file.read()
        file_size = len(data)
    except Exception as e:
        return jsonify({"status": f"Error reading file: {str(e)}"}), 500
    
    return jsonify({
        "status": "success",
        "filename": filename,
        "size": file_size,
        "content_type": file.content_type
    })

if __name__ == "__main__":
    # Enable logging
    import logging
    logging.basicConfig(level=logging.INFO)
    
    # Run the app
    app.run(debug=True, host='0.0.0.0', port=5000)